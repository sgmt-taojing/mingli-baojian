#!/usr/bin/env python3
"""
Desktop-YZX 单步夜间蒸馏（固化脚本模式）
R119 修真：cron 不再做现场决策，改为确定性单步处理。

行为：
1. 读 ~/Desktop/周易-中医/ 目录文件 + 与记账文件比对
2. 选 1 个最旧未处理文件（按优先级 docx > pptx > 小pdf > 大pdf）
3. 提取文本（docx / pptx / pdf 前3页），按主题分块入库
4. FTS5 重建（保留 hit_count 触发器）
5. 输出 JSON 状态
"""
import os, sys, json, hashlib, sqlite3, traceback
from datetime import datetime, timezone, timedelta

ROOT = "/Users/tom/.openclaw-autoclaw/workspace/projects/mingli-baojian"
SRC_DIR = "/Volumes/data1/训练素材-20260816/周易-中医"  # R121：源目录已移至 data1
ACCT_FILE = f"{ROOT}/.openclaw/tmp/yzx-processed.json"
DB_PATH = f"{ROOT}/server/database/yidao.db"
SRC_ID_PREFIX = "SRC-LD-DESKTOP"

# 文件优先级 + 提取后端
EXTRACTORS = [
    (".docx", 100, "docx"),
    (".pptx", 90, "pptx"),
    (".pdf", 50, "pdf"),
]
MIN_BLOCK_LEN = 300
PDF_MAX_PAGES = 3
PDF_MAX_BYTES = 30 * 1024 * 1024  # 30MB 以上的 PDF 跳过（视觉处理范畴）
TIMEOUT_S = 120

def beijing_now():
    return datetime.now(timezone(timedelta(hours=8))).strftime("%Y-%m-%d %H:%M:%S %z").strip()

def fingerprint(text: str) -> str:
    return hashlib.sha1(text.encode("utf-8", errors="ignore")).hexdigest()

def load_acct():
    if not os.path.exists(ACCT_FILE):
        return {}
    try:
        with open(ACCT_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}

def save_acct(d):
    os.makedirs(os.path.dirname(ACCT_FILE), exist_ok=True)
    with open(ACCT_FILE, "w", encoding="utf-8") as f:
        json.dump(d, f, ensure_ascii=False, indent=2)

def select_file(acct: dict):
    if not os.path.isdir(SRC_DIR):
        return None, "src_dir_not_found"
    candidates = []
    for name in os.listdir(SRC_DIR):
        full = os.path.join(SRC_DIR, name)
        if not os.path.isfile(full):
            continue
        if name in acct:
            continue
        if name.startswith("."):
            continue
        ext = os.path.splitext(name)[1].lower()
        size = os.path.getsize(full)
        for ext_pat, prio, kind in EXTRACTORS:
            if ext == ext_pat:
                if kind == "pdf" and size > PDF_MAX_BYTES:
                    continue
                candidates.append((prio, os.path.getmtime(full), full, kind, size))
                break
    if not candidates:
        return None, "no_new_files"
    candidates.sort(key=lambda x: (x[1], -x[0]))  # 旧 mtime 优先；同 mtime 取高优先级
    return candidates[0], None

def extract_text(path: str, kind: str) -> str:
    if kind == "docx":
        from docx import Document
        d = Document(path)
        return "\n".join(p.text for p in d.paragraphs if p.text and p.text.strip())
    if kind == "pptx":
        from pptx import Presentation
        prs = Presentation(path)
        out = []
        for s_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    out.append(f"[slide {s_idx+1}] {shape.text}")
        return "\n".join(out)
    if kind == "pdf":
        import fitz
        doc = fitz.open(path)
        pages = min(PDF_MAX_PAGES, len(doc))
        out = []
        for i in range(pages):
            try:
                out.append(doc[i].get_text())
            except Exception:
                pass
        doc.close()
        return "\n".join(out)
    return ""

def classify_module(title: str, text: str) -> str:
    t = (title + " " + text[:500]).lower()
    if any(k in t for k in ("中医", "望诊", "舌", "脉", "脏腑", "中药", "辨证")):
        return "tcm-wangzhen"
    if any(k in t for k in ("紫微", "斗数")):
        return "ziwei"
    if any(k in t for k in ("风水", "阴宅", "阳宅", "撼龙", "葬书", "龙脉")):
        return "fengshui"
    if any(k in t for k in ("六壬", )) :
        return "liuren"
    if any(k in t for k in ("八字", "四柱", "大运", "流年")):
        return "bazi"
    if any(k in t for k in ("奇门", )) :
        return "qimen"
    if any(k in t for k in ("梅花", "易经", "周易", "六爻", "爻")):
        return "yijing"
    return "yijing"

def chunk_text(text: str, min_len: int = MIN_BLOCK_LEN):
    paras = [p.strip() for p in text.split("\n") if p.strip()]
    blocks, buf = [], ""
    for p in paras:
        if len(buf) + len(p) > 1500:
            if len(buf) >= min_len:
                blocks.append(buf)
            buf = p
        else:
            buf = (buf + "\n" + p).strip()
    if len(buf) >= min_len:
        blocks.append(buf)
    return blocks

def rebuild_fts(conn):
    """FTS5 重建并保留 hit_count 触发器（R111 教训）"""
    cur = conn.cursor()
    # 备份触发器名
    cur.execute("SELECT name FROM sqlite_master WHERE type='trigger' AND tbl_name='kb_formal'")
    triggers = [r[0] for r in cur.fetchall()]
    cur.execute("INSERT INTO kb_formal_fts(kb_formal_fts) VALUES('rebuild')")
    conn.commit()
    return triggers

def ensure_audit_columns(conn):
    """R117/R119 历史经验：kb_formal 缺字段会导致 INSERT 失败。
    这里预检 required 列是否存在，缺则补空列（业务决定，不动 schema）。
    """
    cur = conn.cursor()
    cur.execute("PRAGMA table_info(kb_formal)")
    cols = {r[1] for r in cur.fetchall()}
    required = {"entry_id","module","title","content","src_id","category","keywords","summary",
                "trust_score","version","promoted_at","promoted_from","reviewed_by","hit_count",
                "last_hit","tags","source_ids","confidence","access_level","difficulty","status",
                "created_at","updated_at","audit_status","audit_by","audit_at","audit_notes",
                "model_id","fingerprint","authority"}
    if not required.issubset(cols):
        missing = required - cols
        return False, sorted(missing)
    return True, []

def main():
    out = {"status": "ok", "ts": beijing_now(), "script": "yzx-step-distill.py"}
    acct = load_acct()

    selected, reason = select_file(acct)
    if not selected:
        out["status"] = reason
        print(json.dumps(out, ensure_ascii=False))
        return

    prio, mtime, path, kind, size = selected
    fname = os.path.basename(path)
    out["file"] = fname
    out["size"] = size
    out["kind"] = kind

    try:
        text = extract_text(path, kind)
    except Exception as e:
        out["status"] = "extract_error"
        out["error"] = str(e)[:200]
        # 把这个文件记账为 extract_failed，避免下次重试相同损坏文件
        acct[fname] = {"ts": beijing_now(), "status": "extract_failed", "error": str(e)[:200]}
        save_acct(acct)
        print(json.dumps(out, ensure_ascii=False))
        return

    if len(text) < 200:
        out["status"] = "scanned_or_empty"
        out["text_len"] = len(text)
        acct[fname] = {"ts": beijing_now(), "status": "scanned_or_empty", "text_len": len(text)}
        save_acct(acct)
        print(json.dumps(out, ensure_ascii=False))
        return

    blocks = chunk_text(text)
    if not blocks:
        out["status"] = "no_qualifying_blocks"
        out["text_len"] = len(text)
        acct[fname] = {"ts": beijing_now(), "status": "no_qualifying_blocks"}
        save_acct(acct)
        print(json.dumps(out, ensure_ascii=False))
        return

    conn = sqlite3.connect(DB_PATH, timeout=30)
    ok, missing = ensure_audit_columns(conn)
    if not ok:
        conn.close()
        out["status"] = "schema_mismatch"
        out["missing_columns"] = missing
        print(json.dumps(out, ensure_ascii=False))
        sys.exit(1)

    title_base = os.path.splitext(fname)[0][:60]
    src_id = f"{SRC_ID_PREFIX}-{title_base[:40]}"
    now = beijing_now()
    entry_ids, skipped = [], 0
    try:
        for i, block in enumerate(blocks):
            fp = fingerprint(block)
            entry_id = f"KB-DSK-{fp[:8]}-{i:02d}"
            module = classify_module(title_base, block)
            cur = conn.cursor()
            cur.execute(
                """INSERT OR IGNORE INTO kb_formal
                   (entry_id,module,title,content,src_id,category,keywords,summary,trust_score,
                    version,promoted_at,promoted_from,reviewed_by,hit_count,last_hit,tags,
                    source_ids,confidence,access_level,difficulty,status,created_at,updated_at,
                    audit_status,audit_by,audit_at,audit_notes,model_id,fingerprint,authority)
                   VALUES (?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)""",
                (entry_id, module, f"{title_base}#{i+1}", block[:4000], src_id, "路总流年班",
                 "路总,流年,四墓,桃花,天马", block[:200], 0.85,
                 "v1", now, "yzx-step-distill", "auto-bot", 0, None, "yijing|desktop",
                 src_id, 0.85, "public", "intermediate", "active", now, now,
                 "auto-pending", "yzx-bot", now, "desktop single-step", "yidao-v8",
                 fp, "yijing-desktop")
            )
            if cur.rowcount > 0:
                entry_ids.append(entry_id)
            else:
                skipped += 1
        triggers = rebuild_fts(conn)
        conn.commit()
        conn.close()
        acct[fname] = {"ts": now, "status": "distilled", "blocks": len(blocks),
                       "entry_ids": entry_ids[:5], "total_inserted": len(entry_ids),
                       "skipped_dup": skipped, "triggers": len(triggers)}
        save_acct(acct)
        out["status"] = "distilled"
        out["blocks"] = len(blocks)
        out["inserted"] = len(entry_ids)
        out["skipped_dup"] = skipped
        print(json.dumps(out, ensure_ascii=False))
    except Exception as e:
        conn.rollback()
        conn.close()
        out["status"] = "error"
        out["error"] = str(e)[:200]
        out["trace"] = traceback.format_exc()[:300]
        print(json.dumps(out, ensure_ascii=False))
        sys.exit(1)

if __name__ == "__main__":
    main()